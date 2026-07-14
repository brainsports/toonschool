# -*- coding: utf-8 -*-
"""
툰스쿨에디터 사용설명서 빌드 스크립트
- PPTX(편집 원본, python-pptx) / DOCX(편집 원본, python-docx) / PDF(배포본, Pillow 이미지 기반) + 미리보기 PNG 생성.
- 콘텐츠는 본 파일의 PAGES 공유 데이터 하나로 세 파일이 일치.
- 화면 캡처는 플레이스홀더 프레임(파일명 표시)으로 처리 → 실제 캡처는 PPTX '그림 변경'/DOCX 그림 교체로 대체.
"""
import os
from PIL import Image, ImageDraw, ImageFont

ROOT = os.path.dirname(os.path.abspath(__file__))
BASE = os.path.dirname(ROOT)  # docs/toonschool-editor-manual
ASSETS = os.path.join(ROOT, "_assets")
PREV = os.path.join(BASE, "previews", "page-preview-images")
PDFDIR = os.path.join(BASE, "pdf")
for d in (ASSETS, PREV, PDFDIR):
    os.makedirs(d, exist_ok=True)

FONT_REG = "C:/Windows/Fonts/malgun.ttf"
FONT_BOLD = "C:/Windows/Fonts/malgunbd.ttf"
if not os.path.exists(FONT_BOLD):
    FONT_BOLD = FONT_REG

# ── 디자인 시스템 ──
C_PRIMARY = (0x13, 0x4B, 0x4C)
C_CYAN_BG = (0xDD, 0xEE, 0xEE)
C_WHITE = (0xFF, 0xFF, 0xFF)
C_SOFT = (0xF6, 0xF8, 0xF8)
C_YELLOW = (0xFF, 0xD8, 0x5A)
C_ORANGE = (0xF5, 0x9E, 0x42)
C_RED = (0xD9, 0x4A, 0x4A)
C_BLUE = (0x3B, 0x82, 0xC4)
C_TEXT = (0x22, 0x22, 0x22)
C_SUB = (0x66, 0x66, 0x66)
C_BORDER = (0xD9, 0xD9, 0xD9)
VERSION = "v1.0"
DATE_STR = "2026년 7월"

# A4 가로 mm
W_MM, H_MM = 297, 210
DPI = 150
MM2PX = DPI / 25.4
def mm(v): return int(round(v * MM2PX))
PW, PH = mm(W_MM), mm(H_MM)  # 1754 x 1240

# ── 콘텐츠 (manual-copy.md와 동일) ──
PAGES = [
  {"num":1,"kind":"cover","title":"툰스쿨에디터 사용설명서","desc":"초등학생과 선생님을 위한 만화책 만들기 안내서","shots":["01-cover-finished.png"],"body":["툰스쿨에디터 사용설명서","초등학생과 선생님을 위한 만화책 만들기 안내서",f"{VERSION} · {DATE_STR}"],"prev":"—","next":"툰스쿨에디터 소개"},
  {"num":2,"kind":"intro","title":"툰스쿨에디터 소개","desc":"학교 공부를 내 손으로 만드는 만화책으로 완성해요.","shots":["02-intro-cover.png","02-intro-comic.png","02-intro-story.png","02-intro-quiz.png","02-intro-backcover.png"],"steps":[("1","표지: 단원과 주제를 담은 책 표지를 만들어요."),("2","6컷 만화: 내가 고른 주제로 6장의 만화를 그려요."),("3","세상 속 이야기: 단원을 역사·최신·생활과 연결해요."),("4","OX 퀴즈: 친구들에게 풀어볼 문제 5개를 만들어요."),("5","뒤표지: 지은이와 정보를 담아 책을 완성해요.")],"tip":{"type":"help","title":"알아두세요","text":"한 권의 만화책 안에 표지·만화·이야기·퀴즈·뒤표지가 모두 들어 있어요."},"prev":"표지","next":"시작하기"},
  {"num":3,"kind":"content","title":"시작하기","desc":"로그인하고 에디터로 들어가는 방법이에요.","route":"/login?mode=student → /student/mypage","shots":["03-login.png","03-mypage.png","03-enter-editor.png"],"steps":[("1","로그인 화면에서 아이디와 비밀번호를 넣고 ‘로그인’을 눌러요."),("2","로그인하면 마이페이지가 열려요. 내 작품과 출석을 볼 수 있어요."),("3","위쪽의 ‘툰스쿨 에디터 입장’을 누르면 만화 만들기가 시작돼요.")],"tip":{"type":"help","title":"알아두세요","text":"비밀번호는 다른 사람에게 알려주지 마세요."},"prev":"툰스쿨에디터 소개","next":"전체 제작 과정"},
  {"num":4,"kind":"content","title":"전체 제작 과정","desc":"만화책은 모두 8단계로 만들어져요.","route":"StudentFlowSidebar","shots":["04-flow-sidebar.png"],"steps":[("1","단원 선택 — 학년·과목·단원을 골라요."),("2","주제 만들기 — 핵심어로 만화 주제를 정해요."),("3","대본 만들기 — 6컷 이야기를 만들어요."),("4","표지만들기 — 책 표지를 꾸며요."),("5","만화제작 — 배경·인물·말풍선을 넣어요.")],"tip":{"type":"help","title":"알아두세요","text":"왼쪽 단계 표시에서 완료한 단계에는 체크가 생겨요. 다음 단계는 잠겨 있어요. (6단원정리·7뒤표지·8만화보기 이어짐)"},"prev":"시작하기","next":"학년·학기·과목·단원 선택"},
  {"num":5,"kind":"content","title":"학년·학기·과목·단원 선택","desc":"만들고 싶은 단원을 골라요.","route":"/student/select-unit","shots":["05-unit-step1.png","05-unit-step2.png","05-unit-selected.png"],"steps":[("1","학년과 학기를 골라요. (예: 초4, 1학기)"),("2","‘다음 단계 🚀’를 눌러요."),("3","과목과 대단원, 중단원을 골라요. (예: 과학)"),("4","‘주제 만들기 ✨’를 누르면 다음으로 넘어가요.")],"tip":{"type":"help","title":"알아두세요","text":"선생님이 정해 준 학년·단원만 골라야 할 수도 있어요."},"prev":"전체 제작 과정","next":"핵심어 선택"},
  {"num":6,"kind":"content","title":"핵심어 선택","desc":"만화에 넣을 핵심어를 골라요.","route":"/student/topic (키워드)","shots":["06-keyword-list.png","06-keyword-more.png","06-keyword-selected.png"],"steps":[("1","‘키워드 추천 ✨’를 누르면 핵심어가 나와요."),("2","마음에 드는 핵심어를 눌러 골라요. (최대 4개)"),("3","‘키워드 2개 더 보기 ✨’로 더 볼 수 있어요."),("4","다 골랐으면 ‘이 키워드로 주제 만들기 ✨’를 눌러요.")],"tip":{"type":"caution","title":"주의하세요","text":"핵심어는 최대 4개까지만 고를 수 있어요."},"prev":"학년·학기·과목·단원 선택","next":"만화 주제 만들기"},
  {"num":7,"kind":"content","title":"만화 주제 만들기","desc":"고른 핵심어로 만화 주제를 받아요.","route":"/student/topic (주제)","shots":["07-topic-before.png","07-topic-result.png","07-topic-selected.png"],"steps":[("1","핵심어를 바탕으로 주제가 2개 만들어져요."),("2","‘+2개 더 보기’로 다른 주제를 볼 수 있어요. (최대 10개)"),("3","가장 마음에 드는 주제를 한 개 골라요."),("4","‘대본 만들기 🚀’를 눌러요.")],"tip":{"type":"help","title":"알아두세요","text":"주제가 마음에 없으면 ‘새로운 이야기로 다시 받기 🔄’를 눌러요."},"prev":"핵심어 선택","next":"6컷 대본 만들기"},
  {"num":8,"kind":"content","title":"6컷 대본 만들기","desc":"AI가 6컷 대본을 만들어 줘요.","route":"/student/script","shots":["08-script-before.png","08-script-loading.png","08-script-result-01.png","08-script-result-02.png"],"steps":[("1","‘AI 생성하기’를 눌러요."),("2","대본이 만들어지는 동안 잠시 기다려요. (1·2단계)"),("3","6컷 대본과 핵심 개념 3개, 표지 대화가 만들어져요."),("4","‘표지만들기’를 눌러요.")],"tip":{"type":"help","title":"알아두세요","text":"‘AI로 다시 만들기’를 누르면 대본을 새로 받을 수 있어요."},"prev":"만화 주제 만들기","next":"대본 수정하기"},
  {"num":9,"kind":"content","title":"대본 수정하기","desc":"대본을 내 마음대로 고쳐요.","route":"/student/script","shots":["09-script-cut-edit.png","09-script-key-concept.png","09-script-cover-dialogue.png"],"steps":[("1","‘컷 편집’에서 6컷의 장면과 대사를 고쳐요."),("2","‘핵심 개념’에서 3개의 개념을 확인·수정해요."),("3","‘표지 대화’에서 하나·도윤·서아의 말을 고쳐요."),("4","다 고쳤으면 ‘표지만들기’를 눌러요.")],"tip":{"type":"caution","title":"주의하세요","text":"다음으로 넘가려면 핵심 개념 3개와 표지 대화 3개가 모두 있어야 해요."},"prev":"6컷 대본 만들기","next":"표지 만들기"},
  {"num":10,"kind":"content","title":"표지 만들기","desc":"책 표지를 완성해요.","route":"/student/front-cover","shots":["10-cover-editor.png","10-cover-complete.png"],"steps":[("1","표지 화면이 열리면 과목에 맞는 배경이 깔려 있어요."),("2","‘표지 만들기’를 누르면 제목·핵심 개념·대화가 표지에 들어가요."),("3","표지가 완성되면 ‘만화 만들기’를 눌러요.")],"tip":{"type":"help","title":"알아두세요","text":"표지 배경은 과목마다 자동으로 달라요."},"prev":"대본 수정하기","next":"만화 제작 화면 알아보기"},
  {"num":11,"kind":"content","title":"만화 제작 화면 알아보기","desc":"만화를 만드는 화면을 알아봐요.","route":"/student/comic/full","shots":["11-comic-full.png"],"steps":[("1","왼쪽 세로 도구로 그림을 꾸며요. (선택·캐릭터·대사·말풍선·배경·레이어)"),("2","가운데 2×3 칸에 6컷 만화가 만들어져요."),("3","컷을 누르면 선택되고, 두 번 누르면 크게 볼 수 있어요."),("4","오른쪽 아래 +/- 로 확대·축소하고 ‘화면 맞춤’으로 맞춰요."),("5","만든 내용은 자동으로 저장돼요.")],"tip":{"type":"caution","title":"주의하세요","text":"만화제작은 최대 15분이 걸릴 수 있어요. 배경이 만들어지는 동안 기다려 주세요."},"prev":"표지 만들기","next":"배경과 장면 만들기"},
  {"num":12,"kind":"content","title":"배경과 장면 만들기","desc":"6컷의 배경 그림을 만들어요.","route":"/student/comic/full (배경)","shots":["12-bg-before.png","12-bg-generating.png","12-bg-done.png"],"steps":[("1","‘배경’ 도구를 누르고 컷의 배경 설명을 확인해요."),("2","‘배경 모두 생성’을 누르면 6컷 배경이 차례로 만들어져요."),("3","배경이 다 만들어지면 ‘배경 생성 완료’가 떠요."),("4","설명을 고쳐 ‘수정한 설명으로 다시 만들기’로 다시 만들 수 있어요.")],"tip":{"type":"caution","title":"주의하세요","text":"배경 수정은 컷마다 1번만 할 수 있어요. 사람·글자는 배경에 들어가지 않아요."},"prev":"만화 제작 화면 알아보기","next":"등장인물 넣기"},
  {"num":13,"kind":"content","title":"등장인물 넣기","desc":"만화에 인물을 넣어요.","route":"/student/comic/full (캐릭터)","shots":["13-char-panel.png","13-char-placed.png"],"steps":[("1","‘캐릭터’ 도구를 눌러요."),("2","캐릭터를 골라 캔버스로 끌어다 놓아요."),("3","인물을 눌러 선택하고, 위치와 크기를 조절해요.")],"tip":{"type":"help","title":"알아두세요","text":"여러 인물을 겹쳐 배치할 수 있어요. ‘레이어’ 도구로 순서를 바꿔요."},"prev":"배경과 장면 만들기","next":"말풍선과 글 넣기"},
  {"num":14,"kind":"content","title":"말풍선과 글 넣기","desc":"인물의 대사를 넣어요.","route":"/student/comic/full (말풍선)","shots":["14-bubble-panel.png","14-bubble-added.png"],"steps":[("1","‘말풍선’ 도구를 눌러 말풍선을 골라요."),("2","캔버스에 말풍선을 놓고 대사를 적어요."),("3","‘대사’ 도구에서 대본의 대사를 불러올 수도 있어요."),("4","말풍선의 위치와 크기를 조절해요.")],"tip":{"type":"help","title":"알아두세요","text":"‘대사 생성’을 누르면 대본의 대사가 자동으로 말풍선에 들어가요."},"prev":"등장인물 넣기","next":"6컷 만화 완성하기"},
  {"num":15,"kind":"content","title":"6컷 만화 완성하기","desc":"6컷을 모두 완성하고 다음으로 넘어가요.","route":"/student/comic/full","shots":["15-comic-6cuts.png","15-comic-next.png"],"steps":[("1","6컷에 모두 배경이 들어갔는지 확인해요. (n/6 완료)"),("2","그림이 없는 컷이 있으면 만들어 채워요."),("3","‘단원 정리 →’를 눌러요."),("4","그림이 없는 컷이 남아 있으면 다음으로 넘갈지 물어봐요.")],"tip":{"type":"help","title":"알아두세요","text":"모든 컷이 완성된 상태로 넘어가는 것이 좋아요."},"prev":"말풍선과 글 넣기","next":"세상 속 이야기 만들기"},
  {"num":16,"kind":"content","title":"세상 속 이야기 만들기","desc":"단원을 세상과 연결하는 이야기를 만들어요.","route":"/student/unit-summary","shots":["16-story-before.png","16-story-result.png","16-story-life-tab.png"],"steps":[("1","‘세상 속 이야기 만들기’를 눌러요."),("2","역사 이야기·최신 이야기·생활 연결 탭이 만들어져요."),("3","‘생활 연결’ 탭을 눌러 내 생활과 관련된 이야기를 봐요."),("4","‘세상 속 이야기 완료’를 눌러요.")],"tip":{"type":"help","title":"알아두세요","text":"세 탭 중 하나 이상에 내용이 있으면 완료돼요."},"prev":"6컷 만화 완성하기","next":"OX 퀴즈 만들기"},
  {"num":17,"kind":"content","title":"OX 퀴즈 만들기","desc":"친구들을 위한 OX 문제 5개를 만들어요.","route":"/student/unit-summary","shots":["17-ox-before.png","17-ox-generated.png","17-ox-edit.png"],"steps":[("1","‘OX 문제 만들기’를 누르면 5문제가 만들어져요."),("2","문제를 눌러 내용을 고칠 수 있어요. (‘수정 가능’)"),("3","O 또는 X를 눌러 정답을 정해요."),("4","‘OX 문제 완료’를 누르고 ‘뒤표지 만들기’로 넘어가요.")],"tip":{"type":"help","title":"알아두세요","text":"5문제가 모두 있어야 완료돼요."},"prev":"세상 속 이야기 만들기","next":"뒤표지와 작품 완성"},
  {"num":18,"kind":"content","title":"뒤표지와 작품 완성","desc":"뒤표지를 꾸미고 작품을 완성해요.","route":"/student/back-cover","shots":["18-backcover-form.png","18-backcover-color.png","18-backcover-next.png"],"steps":[("1","지은이·학년/반·과목·단원·주제·만든 날짜를 확인해요."),("2","과목별 배경색을 고르고, 투명도 슬라이더로 조절해요."),("3","‘정보 생성하기’로 정보를 다시 불러올 수 있어요."),("4","‘만화 보기 🖼️’를 누르면 작품이 완성돼요.")],"tip":{"type":"help","title":"알아두세요","text":"‘만화 보기’를 누르면 완성 보상이 주어져요."},"prev":"OX 퀴즈 만들기","next":"작품 보기·공유·문제 해결"},
  {"num":19,"kind":"content","title":"작품 보기·공유·문제 해결","desc":"만화책을 보고 친구에게 자랑해요.","route":"/student/comic/read, /book/:slug","shots":["19-viewer-start.png","19-viewer-flip.png","19-share-modal.png","19-mypage-works.png"],"steps":[("1","‘책 펼치기’를 눌러 만화책을 넘겨 봐요."),("2","‘PDF 다운로드’로 저장하거나 ‘친구에게 자랑하기’로 공유해요."),("3","공유 링크가 만들어지면 ‘링크 복사하기’로 복사해요."),("4","마이페이지에서 내 작품 목록을 볼 수 있어요.")],"tip":{"type":"caution","title":"주의하세요","text":"마이페이지에서 만들다 만 작품을 누르면 이어하기가 제대로 안 될 수 있어요. 그럴 땐 새로 만들기를 시작하세요."},"prev":"뒤표지와 작품 완성","next":"마지막 페이지"},
  {"num":20,"kind":"last","title":"멋진 만화책이 완성되었어요!","desc":"한 권의 만화책을 자랑해 보세요.","shots":["20-final-cover.png","20-final-flip.png"],"body":["여러분이 만든 한 권의 만화책을 자랑해 보세요.",f"툰스쿨에디터 사용설명서 {VERSION} · {DATE_STR}","화면이 바뀌면 이 설명서의 그림과 설명을 바꿔서 다시 만들 수 있어요."],"prev":"작품 보기·공유·문제 해결","next":"—"},
]

def font(size, bold=False):
    return ImageFont.truetype(FONT_BOLD if bold else FONT_REG, size)

def draw_wrap(draw, xy, text, fnt, fill, max_w, line_h):
    x, y = xy
    for raw in text.split("\n"):
        line = ""
        for ch in raw:
            test = line + ch
            if draw.textlength(test, font=fnt) > max_w and line:
                draw.text((x, y), line, font=fnt, fill=fill)
                y += line_h
                line = ch
            else:
                line = test
        draw.text((x, y), line, font=fnt, fill=fill)
        y += line_h
    return y

# ── 이미지 영역 플레이스홀더 PNG (PPTX/DOCX용) ──
def render_image_area(page):
    iw, ih = mm(159), mm(132)
    img = Image.new("RGB", (iw, ih), C_WHITE)
    d = ImageDraw.Draw(img)
    d.rectangle([0, 0, iw-1, ih-1], outline=C_BORDER, width=3)
    d.rectangle([4, 4, iw-5, ih-5], outline=C_CYAN_BG, width=2, fill=C_SOFT)
    # 카메라 아이콘 느낌의 큰 원
    cx, cy = iw//2, ih//2 - 10
    d.ellipse([cx-46, cy-46, cx+46, cy+46], outline=C_PRIMARY, width=5)
    d.ellipse([cx-22, cy-22, cx+22, cy+22], outline=C_PRIMARY, width=4)
    f1 = font(26, True); f2 = font(18); f3 = font(15)
    draw_wrap(d, (20, cy+60), "화면 캡처 필요", f1, C_PRIMARY, iw-40, 32)
    for i, sh in enumerate(page.get("shots", [])):
        draw_wrap(d, (20, cy+100+i*24), f"· {sh}", f3, C_SUB, iw-40, 22)
    return img

# ── 전체 페이지 렌더 (PDF/미리보기용, Pillow) ──
def render_full_page(page):
    img = Image.new("RGB", (PW, PH), C_WHITE)
    d = ImageDraw.Draw(img)
    num = page["num"]; kind = page["kind"]
    if kind == "cover":
        d.rectangle([0, 0, PW, PH], fill=C_PRIMARY)
        d.rectangle([0, PH-90, PW, PH], fill=C_CYAN_BG)
        fb = font(60, True); fs = font(28); fv = font(24)
        draw_wrap(d, (PW//2-460, PH//2-120), page["body"][0], fb, C_WHITE, 920, 76)
        draw_wrap(d, (PW//2-460, PH//2+10), page["body"][1], fs, C_CYAN_BG, 920, 36)
        draw_wrap(d, (PW//2-200, PH-70), page["body"][2], fv, C_PRIMARY, 400, 30)
        return img
    if kind == "last":
        d.rectangle([0, 0, PW, 120], fill=C_PRIMARY)
        d.rectangle([0, PH-90, PW, PH], fill=C_CYAN_BG)
        fb = font(48, True); fs = font(24); fv = font(20)
        y = 180
        for i, line in enumerate(page["body"]):
            f = fb if i == 0 else (fv if i == 2 else fs)
            c = C_WHITE if False else (C_PRIMARY if i == 0 else C_SUB)
            if i == 0:
                draw_wrap(d, (PW//2-560, 40), line, fb, C_WHITE, 1120, 60)
            else:
                y = draw_wrap(d, (PW//2-560, y), line, f, c, 1120, 34)
        # 플레이스홀더 박스
        ix, iy, iw, ih = mm(60), mm(95), mm(177), mm(80)
        d.rectangle([ix, iy, ix+iw, iy+ih], outline=C_BORDER, width=3, fill=C_SOFT)
        draw_wrap(d, (ix+20, iy+ih//2-20), "완성 작품 표지 / 플립북 펼침 화면", font(20, True), C_SUB, iw-40, 28)
        return img
    # content / intro
    # 헤더 바
    d.rectangle([0, 0, PW, mm(36)], fill=C_PRIMARY)
    d.rectangle([0, mm(36), PW, mm(40)], fill=C_YELLOW)
    # STEP 뱃지
    d.rounded_rectangle([mm(16), mm(8), mm(64), mm(30)], radius=6, fill=C_YELLOW)
    draw_wrap(d, (mm(22), mm(11)), f"STEP {num}", font(18, True), C_PRIMARY, 60, 22)
    draw_wrap(d, (mm(72), mm(9)), page["title"], font(24, True), C_WHITE, 1500, 30)
    draw_wrap(d, (mm(72), mm(24)), page["desc"], font(15), C_CYAN_BG, 1500, 20)
    # 왼쪽 이미지 영역
    ix, iy, iw, ih = mm(16), mm(48), mm(159), mm(122)
    d.rectangle([ix, iy, ix+iw, iy+ih], outline=C_BORDER, width=3, fill=C_SOFT)
    cx, cy = ix+iw//2, iy+ih//2-10
    d.ellipse([cx-40, cy-40, cx+40, cy+40], outline=C_PRIMARY, width=5)
    d.ellipse([cx-19, cy-19, cx+19, cy+19], outline=C_PRIMARY, width=4)
    draw_wrap(d, (ix+18, cy+52), "화면 캡처 필요", font(20, True), C_PRIMARY, iw-36, 26)
    for i, sh in enumerate(page.get("shots", [])):
        draw_wrap(d, (ix+18, cy+86+i*22), f"· {sh}", font(13), C_SUB, iw-36, 20)
    # 오른쪽 설명 영역
    rx, ry, rw = mm(183), mm(48), mm(98)
    draw_wrap(d, (rx, ry), "이렇게 해요", font(17, True), C_PRIMARY, rw, 24)
    yy = ry + 32
    for n, txt in page.get("steps", []):
        d.ellipse([rx, yy+2, rx+22, yy+24], fill=C_PRIMARY)
        draw_wrap(d, (rx+5, yy+3), n, font(13, True), C_WHITE, 18, 18)
        yy = draw_wrap(d, (rx+30, yy), txt, font(13), C_TEXT, rw-30, 19) + 4
    # tip 박스
    tip = page.get("tip")
    if tip:
        ty = yy + 6
        col = C_BLUE if tip["type"] == "help" else C_ORANGE
        bg = (0xEA if tip["type"]=="help" else 0xFE, 0xF2 if tip["type"]=="help" else 0xEF, 0xFC if tip["type"]=="help" else 0xE6)
        d.rounded_rectangle([rx, ty, rx+rw, ty+58], radius=6, fill=bg, outline=col, width=2)
        draw_wrap(d, (rx+10, ty+6), tip["title"], font(14, True), col, rw-20, 18)
        draw_wrap(d, (rx+10, ty+26), tip["text"], font(12), C_TEXT, rw-20, 16)
    # 푸터
    d.rectangle([0, mm(190), PW, PH], fill=C_SOFT)
    d.line([0, mm(190), PW, mm(190)], fill=C_BORDER, width=1)
    draw_wrap(d, (mm(16), mm(194)), f"이전: {page['prev']}", font(12), C_SUB, 600, 16)
    draw_wrap(d, (PW//2-60, mm(194)), f"{num} / 20", font(13, True), C_PRIMARY, 120, 16)
    draw_wrap(d, (PW-mm(220), mm(194)), f"{VERSION}  |  다음: {page['next']}", font(12), C_SUB, 400, 16)
    return img

# ── PPTX ──
def build_pptx():
    from pptx import Presentation
    from pptx.util import Mm, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    prs = Presentation()
    prs.slide_width = Mm(297); prs.slide_height = Mm(210)
    blank = prs.slide_layouts[6]
    def rgb(t): return RGBColor(*t)
    def txt(slide, x, y, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tb = slide.shapes.add_textbox(Mm(x), Mm(y), Mm(w), Mm(h))
        tf = tb.text_frame; tf.word_wrap = True
        tf.vertical_anchor = anchor
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = rgb(color)
        r.font.name = "맑은 고딕"
        return tb
    def box(slide, x, y, w, h, fill, line=None, shape=MSO_SHAPE.RECTANGLE):
        sp = slide.shapes.add_shape(shape, Mm(x), Mm(y), Mm(w), Mm(h))
        sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
        if line: sp.line.color.rgb = rgb(line); sp.line.width = Pt(1)
        else: sp.line.fill.background()
        sp.shadow.inherit = False
        return sp
    for pg in PAGES:
        s = prs.slides.add_slide(blank)
        num = pg["num"]; kind = pg["kind"]
        if kind == "cover":
            box(s, 0, 0, 297, 210, C_PRIMARY)
            box(s, 0, 178, 297, 32, C_CYAN_BG)
            txt(s, 40, 70, 217, 40, pg["body"][0], 40, C_WHITE, True, PP_ALIGN.CENTER)
            txt(s, 40, 115, 217, 24, pg["body"][1], 18, C_CYAN_BG, False, PP_ALIGN.CENTER)
            txt(s, 90, 182, 117, 18, pg["body"][2], 16, C_PRIMARY, True, PP_ALIGN.CENTER)
            continue
        if kind == "last":
            box(s, 0, 0, 297, 34, C_PRIMARY)
            box(s, 0, 178, 297, 32, C_CYAN_BG)
            txt(s, 40, 6, 217, 24, pg["body"][0], 24, C_WHITE, True, PP_ALIGN.CENTER)
            txt(s, 40, 80, 217, 20, pg["body"][1], 16, C_SUB, True, PP_ALIGN.CENTER)
            txt(s, 40, 110, 217, 40, pg["body"][2], 14, C_SUB, False, PP_ALIGN.CENTER)
            img_area = render_image_area(pg)
            p = os.path.join(ASSETS, f"img-{num:02d}.png"); img_area.save(p)
            s.shapes.add_picture(p, Mm(60), Mm(95), Mm(177), Mm(80))
            txt(s, 40, 182, 217, 18, f"{VERSION} · {DATE_STR}", 12, C_PRIMARY, True, PP_ALIGN.CENTER)
            continue
        # content / intro
        box(s, 0, 0, 297, 36, C_PRIMARY)
        box(s, 0, 36, 297, 4, C_YELLOW)
        box(s, 16, 8, 48, 22, C_YELLOW, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(s, 18, 10, 44, 18, f"STEP {num}", 14, C_PRIMARY, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
        txt(s, 70, 6, 210, 14, pg["title"], 18, C_WHITE, True)
        txt(s, 70, 22, 210, 12, pg["desc"], 11, C_CYAN_BG)
        # 이미지 영역(플레이스홀더 PNG) — 그림 변경으로 교체 가능
        img_area = render_image_area(pg)
        p = os.path.join(ASSETS, f"img-{num:02d}.png"); img_area.save(p)
        s.shapes.add_picture(p, Mm(16), Mm(48), Mm(159), Mm(122))
        # 오른쪽 설명
        rx = 183
        txt(s, rx, 48, 98, 10, "이렇게 해요", 13, C_PRIMARY, True)
        yy = 64
        for n, t in pg.get("steps", []):
            box(s, rx, yy+1, 7, 7, C_PRIMARY, shape=MSO_SHAPE.OVAL)
            txt(s, rx+1, yy, 7, 8, n, 9, C_WHITE, True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
            txt(s, rx+10, yy, 88, 24, t, 10.5, C_TEXT)
            yy += 22
        tip = pg.get("tip")
        if tip:
            col = C_BLUE if tip["type"] == "help" else C_ORANGE
            bg = (0xEA, 0xF2, 0xFC) if tip["type"] == "help" else (0xFE, 0xEF, 0xE6)
            box(s, rx, yy+4, 98, 30, bg, col, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
            txt(s, rx+5, yy+7, 88, 10, tip["title"], 11, col, True)
            txt(s, rx+5, yy+19, 88, 14, tip["text"], 9.5, C_TEXT)
        # 푸터
        box(s, 0, 190, 297, 20, C_SOFT)
        txt(s, 16, 192, 90, 12, f"이전: {pg['prev']}", 9, C_SUB)
        txt(s, 130, 192, 40, 12, f"{num} / 20", 10, C_PRIMARY, True, PP_ALIGN.CENTER)
        txt(s, 200, 192, 90, 12, f"{VERSION}  |  다음: {pg['next']}", 9, C_SUB)
    out = os.path.join(ROOT, "toonschool-editor-manual-v1.0.pptx")
    prs.save(out); return out

# ── DOCX ──
def build_docx():
    from docx import Document
    from docx.shared import Mm, Pt, RGBColor
    from docx.enum.section import WD_ORIENT
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    doc = Document()
    sec = doc.sections[0]
    sec.page_width = Mm(297); sec.page_height = Mm(210)
    sec.orientation = WD_ORIENT.LANDSCAPE
    for m in ("top_margin","bottom_margin","left_margin","right_margin"):
        setattr(sec, m, Mm(14 if m in ("top_margin","bottom_margin") else 16))
    normal = doc.styles["Normal"]; normal.font.name = "맑은 고딕"; normal.font.size = Pt(12)
    normal.element.rPr.rFonts.set(qn("w:eastAsia"), "맑은 고딕")
    def shade(cell, hexcolor):
        tcPr = cell._tc.get_or_add_tcPr()
        sh = OxmlElement("w:shd"); sh.set(qn("w:val"),"clear"); sh.set(qn("w:fill"), hexcolor)
        tcPr.append(sh)
    def run_fmt(run, size, color, bold=False):
        run.font.size = Pt(size); run.bold = bold; run.font.color.rgb = RGBColor(*color)
        run.font.name = "맑은 고딕"; run.element.rPr.rFonts.set(qn("w:eastAsia"),"맑은 고딕")
    def heading(num, title, desc):
        t = doc.add_table(rows=1, cols=1); t.autofit = True
        c = t.cell(0,0); shade(c, "134B4C")
        p = c.paragraphs[0]; r = p.add_run(f"STEP {num}  {title}"); run_fmt(r, 18, C_WHITE, True)
        p2 = c.add_paragraph(); r2 = p2.add_run(desc); run_fmt(r2, 11, C_CYAN_BG)
    for pg in PAGES:
        num = pg["num"]; kind = pg["kind"]
        if kind == "cover":
            p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            r = p.add_run(pg["body"][0]); run_fmt(r, 32, C_PRIMARY, True)
            p2 = doc.add_paragraph(); p2.alignment = WD_ALIGN_PARAGRAPH.CENTER
            r2 = p2.add_run(pg["body"][1]); run_fmt(r2, 16, C_SUB)
            p3 = doc.add_paragraph(); p3.alignment = WD_ALIGN_PARAGRAPH.CENTER
            r3 = p3.add_run(pg["body"][2]); run_fmt(r3, 14, C_SUB)
        elif kind == "last":
            p = doc.add_paragraph(); p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            r = p.add_run(pg["body"][0]); run_fmt(r, 26, C_PRIMARY, True)
            for line in pg["body"][1:]:
                pp = doc.add_paragraph(); pp.alignment = WD_ALIGN_PARAGRAPH.CENTER
                rr = pp.add_run(line); run_fmt(rr, 13, C_SUB)
        else:
            heading(num, pg["title"], pg["desc"])
            tbl = doc.add_table(rows=1, cols=2); tbl.allow_autofit = False
            tbl.columns[0].width = Mm(165); tbl.columns[1].width = Mm(100)
            lc = tbl.cell(0,0); rc = tbl.cell(0,1)
            lc.width = Mm(165); rc.width = Mm(100)
            img_area = render_image_area(pg)
            p = os.path.join(ASSETS, f"img-{num:02d}.png"); img_area.save(p)
            lp = lc.paragraphs[0]; lp.alignment = WD_ALIGN_PARAGRAPH.CENTER
            lp.add_run().add_picture(p, width=Mm(150))
            hp = rc.paragraphs[0]; hr = hp.add_run("이렇게 해요"); run_fmt(hr, 13, C_PRIMARY, True)
            for n, t in pg.get("steps", []):
                sp = rc.add_paragraph(); sr = sp.add_run(f"{n}. {t}"); run_fmt(sr, 11, C_TEXT)
            tip = pg.get("tip")
            if tip:
                col = C_BLUE if tip["type"]=="help" else C_ORANGE
                tp = rc.add_paragraph(); tr = tp.add_run(f"【{tip['title']}】 {tip['text']}"); run_fmt(tr, 10, col, True)
        # 페이지 나누기 (마지막 페이지 제외)
        if num < 20:
            doc.add_page_break()
    # 꼬리말: 버전
    fpar = sec.footer.paragraphs[0]; fpar.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fr = fpar.add_run(f"툰스쿨에디터 사용설명서 {VERSION} · {DATE_STR}")
    run_fmt(fr, 9, C_SUB)
    out = os.path.join(ROOT, "toonschool-editor-manual-v1.0.docx")
    doc.save(out); return out

# ── PDF + 미리보기 (Pillow) ──
def build_pdf_and_previews():
    pages = []
    for pg in PAGES:
        im = render_full_page(pg)
        fn = os.path.join(PREV, f"page-{pg['num']:02d}.png")
        im.save(fn)
        pages.append(im.convert("RGB"))
    pdf_path = os.path.join(PDFDIR, "toonschool-editor-manual-v1.0.pdf")
    pages[0].save(pdf_path, save_all=True, append_images=pages[1:], resolution=DPI)
    return pdf_path, len(pages)

if __name__ == "__main__":
    pptx_path = build_pptx()
    docx_path = build_docx()
    pdf_path, n = build_pdf_and_previews()
    print("PPTX:", pptx_path)
    print("DOCX:", docx_path)
    print("PDF:", pdf_path, f"({n} pages)")
    print("previews:", PREV)
